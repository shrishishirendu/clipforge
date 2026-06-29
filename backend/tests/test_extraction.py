"""B5: document extraction — deck/summary → key points (FR-10, FR-11)."""
import pytest

from app.services.extraction import DefaultDocumentParser

parser = DefaultDocumentParser()


def test_summary_txt_splits_into_points(tmp_path):
    f = tmp_path / "s.txt"
    f.write_text("- First key point\n\n• Second point\nx\nThird sentence here",
                 encoding="utf-8")
    pts = parser.extract_key_points(str(f), "summary", "txt")
    texts = [p["text"] for p in pts]
    assert "First key point" in texts          # bullet marker stripped
    assert "Second point" in texts
    assert "Third sentence here" in texts
    assert "x" not in texts                     # too short, dropped
    assert all(p["source"] == "Summary" for p in pts)


def test_summary_docx(tmp_path):
    from docx import Document

    doc = Document()
    doc.add_paragraph("Intro point about Acme")
    doc.add_paragraph("")  # blank -> dropped
    doc.add_paragraph("Second important point")
    f = tmp_path / "s.docx"
    doc.save(str(f))

    texts = [p["text"] for p in parser.extract_key_points(str(f), "summary", "docx")]
    assert "Intro point about Acme" in texts
    assert "Second important point" in texts


def test_deck_pptx_one_point_per_slide(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    title_only = prs.slide_layouts[5]
    prs.slides.add_slide(title_only).shapes.title.text = "Welcome to ClipForge"
    prs.slides.add_slide(title_only).shapes.title.text = "The Problem"
    f = tmp_path / "deck.pptx"
    prs.save(str(f))

    pts = parser.extract_key_points(str(f), "deck", "pptx")
    assert pts[0]["source"] == "Slide 1" and "Welcome to ClipForge" in pts[0]["text"]
    assert pts[1]["source"] == "Slide 2" and "Problem" in pts[1]["text"]


def test_unknown_combination_raises():
    with pytest.raises(ValueError):
        parser.extract_key_points("x", "video", "mp4")
