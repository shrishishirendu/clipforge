"""Document extraction — deck + summary → key points (FR-10, FR-11), behind the
DocumentParser interface (arch §10). Default impl uses python-pptx / python-docx /
pypdf. Deck: one key point per slide/page. Summary: one per line/paragraph.
"""
from __future__ import annotations

import re


def _clean(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def _pptx_points(path: str) -> list[dict]:
    from pptx import Presentation

    points = []
    for i, slide in enumerate(Presentation(path).slides, start=1):
        texts = [shape.text_frame.text.strip()
                 for shape in slide.shapes
                 if shape.has_text_frame and shape.text_frame.text.strip()]
        joined = _clean(" ".join(texts))
        if joined:
            points.append({"text": joined, "source": f"Slide {i}"})
    return points


def _pdf_pages(path: str, label: str) -> list[dict]:
    from pypdf import PdfReader

    points = []
    for i, page in enumerate(PdfReader(path).pages, start=1):
        text = _clean(page.extract_text() or "")
        if text:
            points.append({"text": text, "source": f"{label} {i}"})
    return points


def _pdf_text(path: str) -> str:
    from pypdf import PdfReader

    return "\n".join((p.extract_text() or "") for p in PdfReader(path).pages)


def _docx_text(path: str) -> str:
    from docx import Document

    return "\n".join(p.text for p in Document(path).paragraphs)


def _split_summary(text: str) -> list[dict]:
    """One key point per non-trivial line/bullet (FR-11)."""
    points = []
    for raw in text.splitlines():
        line = _clean(raw).lstrip("•-*–·").strip()
        if len(line) >= 3:
            points.append({"text": line, "source": "Summary"})
    return points


class DefaultDocumentParser:
    """python-pptx / python-docx / pypdf based parser."""

    def extract_key_points(self, path: str, asset_type: str, ext: str) -> list[dict]:
        ext = ext.lower()
        if asset_type == "deck":
            if ext == "pptx":
                return _pptx_points(path)
            if ext == "pdf":
                return _pdf_pages(path, "Slide")
        elif asset_type == "summary":
            if ext == "docx":
                return _split_summary(_docx_text(path))
            if ext == "txt":
                with open(path, encoding="utf-8", errors="replace") as fh:
                    return _split_summary(fh.read())
            if ext == "pdf":
                return _split_summary(_pdf_text(path))
        raise ValueError(f"cannot extract key points from {asset_type} '.{ext}'")


def get_document_parser() -> DefaultDocumentParser:
    """Factory for the configured document parser (arch §10)."""
    return DefaultDocumentParser()
